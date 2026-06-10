import io
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from bs4 import BeautifulSoup

def parse_pdf(file_bytes: bytes) -> str:
    pdf_file = io.BytesIO(file_bytes)
    reader = PdfReader(pdf_file)
    text = []
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text.append(page_text)
    return "\n".join(text)

def parse_docx(file_bytes: bytes) -> str:
    docx_file = io.BytesIO(file_bytes)
    try:
        doc = DocxDocument(docx_file)
        text = [p.text for p in doc.paragraphs]
        # Also extract table cell contents if any
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text:
                        text.append(cell.text)
        return "\n".join(text)
    except Exception as e:
        raise ValueError(f"Failed to parse Word Document. Note that legacy binary .doc files are not supported; please convert to .docx. Error: {str(e)}")

def parse_pptx(file_bytes: bytes) -> str:
    pptx_file = io.BytesIO(file_bytes)
    try:
        prs = PptxPresentation(pptx_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        raise ValueError(f"Failed to parse PowerPoint Presentation. Note that legacy binary .ppt files are not supported; please convert to .pptx. Error: {str(e)}")

def parse_html(file_bytes: bytes) -> str:
    try:
        html_content = file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        html_content = file_bytes.decode("latin-1", errors="ignore")
    soup = BeautifulSoup(html_content, "html.parser")
    # Remove script and style elements
    for script in soup(["script", "style"]):
        script.decompose()
    return soup.get_text(separator=" ")

def parse_text(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1", errors="ignore")

def parse_document(file_type: str, file_bytes: bytes) -> str:
    """
    Parses document bytes based on file_type, extracting raw text contents.
    Maps to Apache Tika behavior in original Java project.
    """
    if file_type == "application/pdf":
        return parse_pdf(file_bytes)
    elif file_type in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword"
    ]:
        return parse_docx(file_bytes)
    elif file_type in [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint"
    ]:
        return parse_pptx(file_bytes)
    elif file_type == "text/html":
        return parse_html(file_bytes)
    elif file_type == "text/plain":
        return parse_text(file_bytes)
    else:
        # Fallback to plain text decode if unknown
        return parse_text(file_bytes)
